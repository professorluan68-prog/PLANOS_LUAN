from __future__ import annotations

import logging
import re
from pathlib import Path

from core.lib.classificador import perfil_disciplina
from core.normalizacao import normalizar as normalizar_texto

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - fallback seguro quando dependencia nao estiver instalada
    Presentation = None


logger = logging.getLogger(__name__)

_ROTULOS_METADADOS = {
    "lingua portuguesa",
    "ensino medio",
    "1 serie",
    "1a serie",
    "3 bimestre",
    "3o bimestre",
    "material digital",
    "conteudos",
    "conteudo",
    "objetivos",
    "objetivos da aula",
}

_MARCADORES_BLOCOS = {
    "para comecar": "Para comecar",
    "foco no conteudo": "Foco no conteudo",
    "na pratica": "Na pratica",
    "encerramento": "Encerramento",
    "pause e responda": "Na pratica",
    "com suas palavras": "Para comecar",
    "virem e conversem": "Para comecar",
    "todo mundo escreve": "Na pratica",
    "hora da leitura": "Foco no conteudo",
    "de olho no modelo": "Foco no conteudo",
}

_TECNICAS_CONHECIDAS = [
    "Virem e conversem",
    "Todo mundo escreve",
    "Com suas palavras",
    "Hora da leitura",
    "De olho no modelo",
    "Pause e responda",
]


def _limpar_texto(texto: str) -> str:
    return re.sub(r"\s+", " ", str(texto or "")).strip(" -:\t\r\n")


def _segmentar_texto_shape(texto: str) -> list[str]:
    segmentos = []
    for parte in re.split(r"[\r\n]+", str(texto or "")):
        parte = _limpar_texto(parte)
        if parte:
            segmentos.append(parte)
    return segmentos


def _normalizar_compacto(texto: str) -> str:
    return normalizar_texto(texto).replace(" ", "")


def _extrair_numero_aula(texto: str) -> str:
    match = re.search(r"\baula\s*(?:n[.o]?\s*)?(\d{1,3})\b", normalizar_texto(texto))
    return match.group(1) if match else ""


def _texto_eh_generico(texto: str, nome_arquivo: str = "") -> bool:
    texto_limpo = _limpar_texto(texto)
    if not texto_limpo:
        return True

    base = normalizar_texto(texto_limpo)
    if len(texto_limpo) < 6:
        return True
    if base in _ROTULOS_METADADOS:
        return True
    if re.fullmatch(r"aula\s*\d{1,3}", base):
        return True
    if re.fullmatch(r"\d+\s*(ano|serie)", base):
        return True
    if re.fullmatch(r"\d+\s*bimestre", base):
        return True
    if base == _normalizar_compacto(nome_arquivo) or base == normalizar_texto(nome_arquivo):
        return True
    if any(token in base for token in ["arquivo", ".pptx", ".pdf"]):
        return True
    return False


def titulo_pptx_valido(titulo: str, nome_arquivo: str = "") -> bool:
    return not _texto_eh_generico(titulo, nome_arquivo=nome_arquivo)


def _coletar_textos_slide(slide) -> list[str]:
    textos = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for trecho in _segmentar_texto_shape(getattr(shape, "text", "")):
            textos.append(trecho)
    return textos


def _extrair_titulo(slides: list[dict], nome_arquivo: str) -> str:
    candidatos: list[tuple[int, str]] = []

    for slide in slides[:3]:
        for indice, texto in enumerate(slide.get("textos", [])):
            texto_limpo = _limpar_texto(texto)
            if not titulo_pptx_valido(texto_limpo, nome_arquivo=nome_arquivo):
                continue
            base = normalizar_texto(texto_limpo)
            if any(rotulo in base for rotulo in ["objetivos", "conteudos", "encerramento"]):
                continue

            pontuacao = 0
            if slide.get("numero") == 1:
                pontuacao += 60
            if indice <= 2:
                pontuacao += 20
            pontuacao += min(len(texto_limpo), 80)
            if ":" in texto_limpo:
                pontuacao -= 10
            if re.search(r"\bparte\s+\d+\b", base):
                pontuacao += 5
            candidatos.append((pontuacao, texto_limpo))

    if candidatos:
        candidatos.sort(key=lambda item: item[0], reverse=True)
        return candidatos[0][1][:160]
    return ""


def _extrair_itens_por_rotulo(slides: list[dict], rotulos: tuple[str, ...]) -> list[str]:
    rotulos_norm = {normalizar_texto(rotulo) for rotulo in rotulos}
    itens: list[str] = []
    capturando = False

    for slide in slides[:5]:
        for texto in slide.get("textos", []):
            texto_limpo = _limpar_texto(texto)
            texto_norm = normalizar_texto(texto_limpo)
            if texto_norm in rotulos_norm:
                capturando = True
                continue
            if capturando and texto_norm in _ROTULOS_METADADOS:
                break
            if capturando and titulo_pptx_valido(texto_limpo):
                itens.append(texto_limpo)
        if capturando and itens:
            break

    itens_limpos = []
    vistos = set()
    for item in itens:
        base = normalizar_texto(item)
        if base and base not in vistos:
            itens_limpos.append(item)
            vistos.add(base)
    return itens_limpos[:8]


def _extrair_blocos(slides: list[dict]) -> dict[str, str]:
    blocos = {
        "Para comecar": "",
        "Foco no conteudo": "",
        "Na pratica": "",
        "Encerramento": "",
    }
    bloco_atual = ""

    for slide in slides:
        for texto in slide.get("textos", []):
            texto_limpo = _limpar_texto(texto)
            texto_norm = normalizar_texto(texto_limpo)
            if texto_norm in _MARCADORES_BLOCOS:
                bloco_atual = _MARCADORES_BLOCOS[texto_norm]
                continue
            if bloco_atual and texto_limpo and titulo_pptx_valido(texto_limpo):
                bloco_existente = blocos.get(bloco_atual, "")
                blocos[bloco_atual] = f"{bloco_existente} {texto_limpo}".strip()

    return {chave: valor.strip() for chave, valor in blocos.items() if valor.strip()}


def _detectar_tecnicas(texto_completo: str) -> list[str]:
    base = normalizar_texto(texto_completo)
    tecnicas = []
    for tecnica in _TECNICAS_CONHECIDAS:
        if normalizar_texto(tecnica) in base:
            tecnicas.append(tecnica)
    return tecnicas


def _montar_texto_completo(
    titulo: str,
    disciplina: str,
    serie: str,
    bimestre: str,
    aula: str,
    conteudos: list[str],
    objetivos: list[str],
    blocos: dict[str, str],
    slides: list[dict],
) -> str:
    partes = []
    for item in [disciplina, serie, bimestre, aula, titulo]:
        if _limpar_texto(item):
            partes.append(_limpar_texto(item))

    if objetivos:
        partes.append("Objetivos da aula")
        partes.extend(objetivos)

    if conteudos:
        partes.append("Conteudos")
        partes.extend(conteudos)

    for marcador in ("Para comecar", "Foco no conteudo", "Na pratica", "Encerramento"):
        if blocos.get(marcador):
            partes.append(marcador)
            partes.append(blocos[marcador])

    for slide in slides:
        partes.append(f"Slide {slide['numero']}")
        partes.extend(slide.get("textos", []))

    return "\n".join(parte for parte in partes if _limpar_texto(parte))


def extrair_estrutura_pptx(caminho_pptx: str) -> dict:
    if Presentation is None:
        raise RuntimeError("Dependencia python-pptx nao instalada.")

    presentation = Presentation(caminho_pptx)
    slides = []
    for numero, slide in enumerate(presentation.slides, start=1):
        textos = _coletar_textos_slide(slide)
        slides.append(
            {
                "numero": numero,
                "titulo": textos[0] if textos else "",
                "textos": textos,
            }
        )

    nome_arquivo = Path(caminho_pptx).stem
    titulo = _extrair_titulo(slides, nome_arquivo=nome_arquivo)
    conteudos = _extrair_itens_por_rotulo(slides, ("Conteudos", "Conteudo"))
    objetivos = _extrair_itens_por_rotulo(slides, ("Objetivos", "Objetivos da aula"))
    blocos = _extrair_blocos(slides)

    todos_os_textos = [texto for slide in slides for texto in slide.get("textos", [])]
    texto_consolidado = "\n".join(todos_os_textos)
    disciplina = next((texto for texto in todos_os_textos[:12] if "portugues" in normalizar_texto(texto)), "Lingua Portuguesa")
    serie = next((texto for texto in todos_os_textos[:16] if re.search(r"\b1\s*(?:a\s*)?(?:serie|ano)\b", normalizar_texto(texto))), "1a Serie")
    bimestre = next((texto for texto in todos_os_textos[:16] if "bimestre" in normalizar_texto(texto)), "")
    aula_numero = _extrair_numero_aula(texto_consolidado) or _extrair_numero_aula(nome_arquivo)
    aula = f"Aula {aula_numero}" if aula_numero else ""

    if not titulo and conteudos:
        titulo = conteudos[0]

    texto_completo = _montar_texto_completo(
        titulo=titulo,
        disciplina=disciplina,
        serie=serie,
        bimestre=bimestre,
        aula=aula,
        conteudos=conteudos,
        objetivos=objetivos,
        blocos=blocos,
        slides=slides,
    )
    tecnicas = _detectar_tecnicas(texto_completo)

    logger.info("[EXTRACAO] Titulo extraido: %s", titulo)
    logger.info("[EXTRACAO] Conteudos extraidos: %s", conteudos)
    logger.info("[EXTRACAO] Objetivos extraidos: %s", objetivos)
    logger.info("[EXTRACAO] Tecnicas detectadas: %s", tecnicas)
    logger.info("[EXTRACAO] Blocos detectados: %s", list(blocos.keys()))

    return {
        "fonte": "pptx",
        "titulo": titulo,
        "disciplina": disciplina,
        "serie": serie,
        "bimestre": bimestre,
        "aula": aula,
        "conteudos": conteudos,
        "objetivos": objetivos,
        "slides": slides,
        "blocos": blocos,
        "tecnicas_detectadas": tecnicas,
        "texto_completo": texto_completo,
    }


def estrutura_pptx_para_dados_aula(estrutura: dict) -> dict:
    titulo = _limpar_texto(estrutura.get("titulo", ""))
    conteudos = [item for item in estrutura.get("conteudos", []) if _limpar_texto(item)]
    objetivos = [item for item in estrutura.get("objetivos", []) if _limpar_texto(item)]
    aula = _limpar_texto(estrutura.get("aula", ""))
    tema = titulo if titulo_pptx_valido(titulo) else (conteudos[0] if conteudos else aula)
    material = f"{aula} - {tema}".strip(" -") if aula else tema
    aprendizagem = "; ".join(objetivos[:3])

    return {
        "tema": tema,
        "material": material,
        "aprendizagem": aprendizagem,
        "texto_base": estrutura.get("texto_completo", ""),
        "conteudos": conteudos,
        "objetivos": objetivos,
        "blocos_pedagogicos": estrutura.get("blocos", {}) or {},
        "fonte_extracao": "pptx",
    }


def eh_cenario_piloto_pptx(disciplina: str, turma: str) -> bool:
    if perfil_disciplina(disciplina, turma=turma) != "lingua_portuguesa_em":
        return False
    turma_norm = normalizar_texto(turma)
    return bool(re.search(r"(?<!\d)1\s*(?:a\s*)?(?:serie|ano)\b", turma_norm))


def encontrar_pptx_correspondente(caminho_pdf: str, disciplina: str, turma: str) -> str | None:
    if not eh_cenario_piloto_pptx(disciplina, turma):
        return None

    pasta = Path(caminho_pdf).parent
    if not pasta.exists():
        return None

    candidatos = [
        caminho for caminho in pasta.glob("*.pptx")
        if not caminho.name.startswith("~$")
    ]
    if not candidatos:
        return None

    numero_pdf = _extrair_numero_aula(Path(caminho_pdf).stem)
    if len(candidatos) == 1:
        unico = str(candidatos[0])
        logger.info("[EXTRACAO] PPTX correspondente encontrado: %s", unico)
        return unico

    melhor_caminho = None
    melhor_pontuacao = -1

    for candidato in candidatos:
        pontuacao = 0
        nome_norm = normalizar_texto(candidato.stem)
        if numero_pdf and re.search(rf"\baula\s*0*{re.escape(numero_pdf)}\b", nome_norm):
            pontuacao += 50
        elif numero_pdf and re.search(rf"(?<!\d)0*{re.escape(numero_pdf)}(?!\d)", nome_norm):
            pontuacao += 20

        try:
            estrutura = extrair_estrutura_pptx(str(candidato))
            aula_estrutura = _extrair_numero_aula(estrutura.get("aula", "") or estrutura.get("texto_completo", ""))
            if numero_pdf and aula_estrutura == numero_pdf:
                pontuacao += 100
        except Exception:
            continue

        if pontuacao > melhor_pontuacao:
            melhor_pontuacao = pontuacao
            melhor_caminho = str(candidato)

    if melhor_caminho and melhor_pontuacao > 0:
        logger.info("[EXTRACAO] PPTX correspondente encontrado: %s", melhor_caminho)
        return melhor_caminho
    return None
